from pathlib import Path, PurePosixPath

import pytest
from fsspec.implementations.http import HTTPFileSystem
from fsspec.implementations.local import LocalFileSystem
from gcsfs import GCSFileSystem
from kedro.io.core import PROTOCOL_DELIMITER, DatasetError, Version
from pptx import Presentation
from s3fs.core import S3FileSystem

from kedro_datasets.openxml import PptxDataset


@pytest.fixture
def filepath_pptx(tmp_path):
    return (tmp_path / "test.pptx").as_posix()


@pytest.fixture
def pptx_dataset(filepath_pptx, fs_args):
    return PptxDataset(filepath=filepath_pptx, fs_args=fs_args)


@pytest.fixture
def versioned_pptx_dataset(filepath_pptx, load_version, save_version):
    return PptxDataset(
        filepath=filepath_pptx, version=Version(load_version, save_version)
    )


@pytest.fixture
def dummy_data() -> Presentation:
    test_pres = Presentation()
    slide = test_pres.slides.add_slide(test_pres.slide_layouts[0])
    title = slide.shapes.title
    title.text = "Hello, this is dummy data pptx presentation."
    return test_pres


class TestPptxDataset:
    def test_save_and_load(self, pptx_dataset, dummy_data):
        """Test saving and reloading the dataset."""
        pptx_dataset.save(dummy_data)
        reloaded = pptx_dataset.load()
        assert (
            dummy_data.slides[0].shapes.title.text
            == reloaded.slides[0].shapes.title.text
        )
        assert pptx_dataset._fs_open_args_load == {}
        assert pptx_dataset._fs_open_args_save == {"mode": "wb"}

    def test_exists(self, pptx_dataset, dummy_data):
        """Test `exists` method invocation for both existing and
        nonexistent dataset."""
        assert not pptx_dataset.exists()
        pptx_dataset.save(dummy_data)
        assert pptx_dataset.exists()

    @pytest.mark.parametrize(
        "fs_args",
        [{"open_args_load": {"mode": "rb", "compression": "gzip"}}],
        indirect=True,
    )
    def test_open_extra_args(self, pptx_dataset, fs_args):
        assert pptx_dataset._fs_open_args_load == fs_args["open_args_load"]
        assert pptx_dataset._fs_open_args_save == {"mode": "wb"}  # default unchanged

    def test_load_missing_file(self, pptx_dataset):
        """Check the error when trying to load missing file."""
        pattern = r"Failed while loading data from dataset kedro_datasets.openxml.pptx_dataset.PptxDataset\(.*\)"
        with pytest.raises(DatasetError, match=pattern):
            pptx_dataset.load()

    @pytest.mark.parametrize(
        "filepath,instance_type",
        [
            ("s3://bucket/file.pptx", S3FileSystem),
            ("file:///tmp/test.pptx", LocalFileSystem),
            ("/tmp/test.pptx", LocalFileSystem),
            ("gcs://bucket/file.pptx", GCSFileSystem),
            ("https://example.com/file.pptx", HTTPFileSystem),
        ],
    )
    def test_protocol_usage(self, filepath, instance_type):
        dataset = PptxDataset(filepath=filepath)
        assert isinstance(dataset._fs, instance_type)

        path = filepath.split(PROTOCOL_DELIMITER, 1)[-1]

        assert str(dataset._filepath) == path
        assert isinstance(dataset._filepath, PurePosixPath)

    def test_catalog_release(self, mocker):
        fs_mock = mocker.patch("fsspec.filesystem").return_value
        filepath = "test.pptx"
        dataset = PptxDataset(filepath=filepath)
        dataset.release()
        fs_mock.invalidate_cache.assert_called_once_with(filepath)


class TestPptxDatasetVersioned:
    def test_version_str_repr(self, load_version, save_version):
        """Test that version is in string representation of the class instance
        when applicable."""
        filepath = "test.pptx"
        ds = PptxDataset(filepath=filepath)
        ds_versioned = PptxDataset(
            filepath=filepath, version=Version(load_version, save_version)
        )
        assert filepath in str(ds)
        assert "version" not in str(ds)

        assert filepath in str(ds_versioned)
        ver_str = f"version=Version(load={load_version}, save='{save_version}')"
        assert ver_str in str(ds_versioned)
        assert "PptxDataset" in str(ds_versioned)
        assert "PptxDataset" in str(ds)
        assert "protocol" in str(ds_versioned)
        assert "protocol" in str(ds)

    def test_save_and_load(self, versioned_pptx_dataset, dummy_data):
        """Test that saved and reloaded data matches the original one for
        the versioned dataset."""
        versioned_pptx_dataset.save(dummy_data)
        reloaded = versioned_pptx_dataset.load()
        assert (
            dummy_data.slides[0].shapes.title.text
            == reloaded.slides[0].shapes.title.text
        )

    def test_no_versions(self, versioned_pptx_dataset):
        """Check the error if no versions are available for load."""
        pattern = r"Did not find any versions for kedro_datasets.openxml.pptx_dataset.PptxDataset\(.+\)"
        with pytest.raises(DatasetError, match=pattern):
            versioned_pptx_dataset.load()

    def test_exists(self, versioned_pptx_dataset, dummy_data):
        """Test `exists` method invocation for versioned dataset."""
        assert not versioned_pptx_dataset.exists()
        versioned_pptx_dataset.save(dummy_data)
        assert versioned_pptx_dataset.exists()

    def test_prevent_overwrite(self, versioned_pptx_dataset, dummy_data):
        """Check the error when attempting to override the dataset if the
        corresponding pptx file for a given save version already exists."""
        versioned_pptx_dataset.save(dummy_data)
        pattern = (
            r"Save path \'.+\' for kedro_datasets.openxml.pptx_dataset.PptxDataset\(.+\) must "
            r"not exist if versioning is enabled\."
        )
        with pytest.raises(DatasetError, match=pattern):
            versioned_pptx_dataset.save(dummy_data)

    @pytest.mark.parametrize(
        "load_version", ["2019-01-01T23.59.59.999Z"], indirect=True
    )
    @pytest.mark.parametrize(
        "save_version", ["2019-01-02T00.00.00.000Z"], indirect=True
    )
    def test_save_version_warning(
        self, versioned_pptx_dataset, load_version, save_version, dummy_data
    ):
        """Check the warning when saving to the path that differs from
        the subsequent load path."""
        pattern = (
            rf"Save version '{save_version}' did not match load version "
            rf"'{load_version}' for kedro_datasets.openxml.pptx_dataset.PptxDataset\(.+\)"
        )
        with pytest.warns(UserWarning, match=pattern):
            versioned_pptx_dataset.save(dummy_data)

    def test_http_filesystem_no_versioning(self):
        pattern = "Versioning is not supported for HTTP protocols."

        with pytest.raises(DatasetError, match=pattern):
            PptxDataset(
                filepath="https://example.com/file.pptx", version=Version(None, None)
            )

    def test_versioning_existing_dataset(
        self, pptx_dataset, versioned_pptx_dataset, dummy_data
    ):
        """Check the error when attempting to save a versioned dataset on top of an
        already existing (non-versioned) dataset."""
        pptx_dataset.save(dummy_data)
        assert pptx_dataset.exists()
        assert pptx_dataset._filepath == versioned_pptx_dataset._filepath
        pattern = (
            f"(?=.*file with the same name already exists in the directory)"
            f"(?=.*{versioned_pptx_dataset._filepath.parent.as_posix()})"
        )
        with pytest.raises(DatasetError, match=pattern):
            versioned_pptx_dataset.save(dummy_data)

        # Remove non-versioned dataset and try again
        Path(pptx_dataset._filepath.as_posix()).unlink()
        versioned_pptx_dataset.save(dummy_data)
        assert versioned_pptx_dataset.exists()
